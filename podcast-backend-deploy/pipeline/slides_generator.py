"""
Slides Generator Module for Slidev PPT Generation.

This module provides functions to convert podcast scripts into Slidev-formatted
Markdown presentations using LLM-based key points extraction.

**Feature: slidev-ppt-generator**
"""
import re
from typing import Dict, Any, Optional


def extract_key_points(
    cfg: Dict[str, Any],
    script: str,
    title: str,
    style: str = "professional"
) -> str:
    """
    使用 LLM 从播客脚本中提取关键要点，生成 Slidev Markdown
    
    Args:
        cfg: 配置信息，包含 Hunyuan API 凭证
        script: 播客脚本文本
        title: 播客标题
        style: 幻灯片风格 (professional/minimal/creative)
    
    Returns:
        Slidev 格式的 Markdown 字符串
    
    Raises:
        ValueError: 如果脚本为空
        RuntimeError: 如果 LLM 调用失败
    """
    if not script or not script.strip():
        raise ValueError("脚本内容不能为空")
    
    from clients.hunyuan_api_client import HunyuanAPIClient
    
    # Initialize Hunyuan API client
    api = HunyuanAPIClient(
        secret_id=cfg.get("hunyuan_api_secret_id", ""),
        secret_key=cfg.get("hunyuan_api_secret_key", ""),
        region=cfg.get("hunyuan_api_region", "ap-beijing"),
        model=cfg.get("hunyuan_api_model", "hunyuan-turbos-latest"),
        temperature=cfg.get("hunyuan_api_temperature", 0.8),
        top_p=cfg.get("hunyuan_api_top_p", 0.8),
        max_tokens=cfg.get("hunyuan_api_max_tokens", 8000),
    )
    
    # Style-specific theme mapping
    theme_map = {
        "professional": "seriph",
        "minimal": "default",
        "creative": "apple-basic"
    }
    theme = theme_map.get(style, "seriph")

    # Build the LLM prompt for key points extraction
    prompt = f"""你是一位专业的演示文稿设计师。请将以下播客脚本转换为 Slidev 格式的 Markdown 演示文稿。

【重要要求】
1. 提取关键要点：从脚本中提取核心观点、数据、结论，而非逐字复制原文
2. 结构化呈现：将内容组织成清晰的幻灯片结构
3. 精简内容：每张幻灯片的文字要精炼，适合演示

【输出格式要求】
1. 必须以 YAML frontmatter 开头，包含 theme: {theme}
2. 使用 --- 分隔每张幻灯片
3. 第一张必须是标题幻灯片，标题为：{title}
4. 最后一张必须是总结/结论幻灯片（包含"总结"、"结论"、"Summary"或"Conclusion"关键词）
5. 中间至少包含 3 张内容幻灯片
6. 根据内容类型使用不同布局：
   - 默认布局：普通内容
   - layout: two-cols：对比或并列内容
   - layout: center：重要引用或强调
   - layout: image-right：配图说明（使用占位符）

【播客脚本】
{script[:8000]}

【输出示例格式】
---
theme: {theme}
title: {title}
class: text-center
transition: slide-left
---

# {title}

基于 AI 播客内容生成

---

# 核心要点

- 要点 1
- 要点 2
- 要点 3

---
layout: two-cols
---

# 详细分析

左侧内容...

::right::

右侧内容...

---
layout: center
---

# 总结

关键结论和行动建议

请直接输出 Slidev Markdown 内容，不要包含任何解释或说明。"""

    messages = [
        {"Role": "system", "Content": "你是一个专业的演示文稿设计师，擅长将长文本内容提炼为精简的幻灯片要点。"},
        {"Role": "user", "Content": prompt},
    ]
    
    try:
        resp = api.chat(messages, stream=False)
        
        # Extract content from response
        content = ""
        choices = resp.get("Choices") or resp.get("choices") or []
        if choices:
            msg = choices[0].get("Message") or choices[0].get("message") or {}
            content = msg.get("Content") or msg.get("content") or ""
        
        if not content:
            raise RuntimeError("LLM 返回内容为空")
        
        # Validate and clean the output
        markdown = _validate_and_clean_slidev_markdown(content, title, theme)
        
        return markdown
        
    except Exception as e:
        raise RuntimeError(f"生成幻灯片失败: {str(e)}")



def _validate_and_clean_slidev_markdown(content: str, title: str, theme: str) -> str:
    """
    验证并清理 LLM 生成的 Slidev Markdown 内容
    
    Args:
        content: LLM 生成的原始内容
        title: 演示文稿标题
        theme: Slidev 主题
    
    Returns:
        清理后的有效 Slidev Markdown
    """
    # Remove any markdown code block wrappers
    content = re.sub(r'^```(?:markdown|md)?\s*\n?', '', content.strip())
    content = re.sub(r'\n?```\s*$', '', content)
    
    # Ensure frontmatter exists
    if not content.startswith('---'):
        frontmatter = f"""---
theme: {theme}
title: {title}
class: text-center
transition: slide-left
---

"""
        content = frontmatter + content
    else:
        # Ensure theme is in frontmatter
        frontmatter_match = re.match(r'^---\s*\n(.*?)\n---', content, re.DOTALL)
        if frontmatter_match:
            frontmatter_content = frontmatter_match.group(1)
            if 'theme:' not in frontmatter_content:
                # Add theme to frontmatter
                new_frontmatter = f"theme: {theme}\n" + frontmatter_content
                content = f"---\n{new_frontmatter}\n---" + content[frontmatter_match.end():]
    
    return content.strip()


def render_preview_html(markdown: str) -> str:
    """
    将 Slidev Markdown 渲染为静态 HTML 预览
    
    Args:
        markdown: Slidev 格式的 Markdown
    
    Returns:
        HTML 字符串
    """
    if not markdown or not markdown.strip():
        return "<div class='error'>Markdown 内容为空</div>"
    
    # Parse slides from markdown
    slides = parse_slidev_markdown(markdown)
    
    # Generate HTML for each slide
    html_slides = []
    for i, slide in enumerate(slides):
        slide_html = _render_slide_to_html(slide, i)
        html_slides.append(slide_html)
    
    # Wrap in container with basic styling
    html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Slides Preview</title>
    <style>
        * {{ box-sizing: border-box; margin: 0; padding: 0; }}
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; background: #1a1a2e; }}
        .slides-container {{ max-width: 960px; margin: 0 auto; padding: 20px; }}
        .slide {{ background: white; border-radius: 8px; padding: 40px; margin-bottom: 20px; min-height: 400px; box-shadow: 0 4px 6px rgba(0,0,0,0.3); }}
        .slide h1 {{ font-size: 2em; margin-bottom: 20px; color: #333; }}
        .slide h2 {{ font-size: 1.5em; margin-bottom: 15px; color: #444; }}
        .slide ul {{ margin-left: 20px; line-height: 1.8; }}
        .slide li {{ margin-bottom: 8px; }}
        .slide p {{ line-height: 1.6; margin-bottom: 10px; }}
        .slide.center {{ display: flex; flex-direction: column; justify-content: center; align-items: center; text-align: center; }}
        .slide.two-cols {{ display: grid; grid-template-columns: 1fr 1fr; gap: 20px; }}
        .slide-number {{ position: absolute; bottom: 10px; right: 15px; color: #999; font-size: 0.8em; }}
        .slide-wrapper {{ position: relative; }}
        .nav-hint {{ text-align: center; color: #666; padding: 20px; }}
    </style>
</head>
<body>
    <div class="slides-container">
        {''.join(html_slides)}
    </div>
    <div class="nav-hint">共 {len(slides)} 张幻灯片</div>
</body>
</html>"""
    
    return html



def parse_slidev_markdown(markdown: str) -> list:
    """
    解析 Slidev Markdown 为幻灯片列表
    
    Args:
        markdown: Slidev 格式的 Markdown
    
    Returns:
        幻灯片内容列表，每个元素包含 layout 和 content
    """
    if not markdown:
        return []
    
    # Split by slide separator (---)
    # First, handle frontmatter
    parts = re.split(r'\n---\s*\n', markdown)
    
    slides = []
    for i, part in enumerate(parts):
        part = part.strip()
        if not part:
            continue
        
        # Skip frontmatter (first part starting with ---)
        if i == 0 and part.startswith('---'):
            # Extract frontmatter content
            fm_match = re.match(r'^---\s*\n(.*)', part, re.DOTALL)
            if fm_match:
                part = fm_match.group(1).strip()
                # Check if this is just frontmatter
                if not any(line.startswith('#') for line in part.split('\n')):
                    continue
        
        # Parse layout from frontmatter-like block at start of slide
        layout = "default"
        content = part
        
        layout_match = re.match(r'^layout:\s*(\S+)\s*\n', part)
        if layout_match:
            layout = layout_match.group(1)
            content = part[layout_match.end():].strip()
        
        slides.append({
            "layout": layout,
            "content": content
        })
    
    return slides


def _render_slide_to_html(slide: dict, index: int) -> str:
    """
    将单张幻灯片渲染为 HTML
    
    Args:
        slide: 幻灯片数据 (layout, content)
        index: 幻灯片索引
    
    Returns:
        HTML 字符串
    """
    layout = slide.get("layout", "default")
    content = slide.get("content", "")
    
    # Convert markdown to basic HTML
    html_content = _markdown_to_html(content)
    
    # Apply layout class
    layout_class = ""
    if layout == "center":
        layout_class = "center"
    elif layout == "two-cols":
        layout_class = "two-cols"
    
    return f"""<div class="slide-wrapper">
    <div class="slide {layout_class}">
        {html_content}
        <span class="slide-number">{index + 1}</span>
    </div>
</div>
"""


def _markdown_to_html(md: str) -> str:
    """
    简单的 Markdown 到 HTML 转换
    
    Args:
        md: Markdown 文本
    
    Returns:
        HTML 字符串
    """
    if not md:
        return ""
    
    lines = md.split('\n')
    html_lines = []
    in_list = False
    
    for line in lines:
        line = line.rstrip()
        
        # Skip layout declarations and slot markers
        if line.startswith('layout:') or line.startswith('::'):
            continue
        
        # Headers
        if line.startswith('# '):
            if in_list:
                html_lines.append('</ul>')
                in_list = False
            html_lines.append(f'<h1>{line[2:]}</h1>')
        elif line.startswith('## '):
            if in_list:
                html_lines.append('</ul>')
                in_list = False
            html_lines.append(f'<h2>{line[3:]}</h2>')
        elif line.startswith('### '):
            if in_list:
                html_lines.append('</ul>')
                in_list = False
            html_lines.append(f'<h3>{line[4:]}</h3>')
        # List items
        elif line.startswith('- ') or line.startswith('* '):
            if not in_list:
                html_lines.append('<ul>')
                in_list = True
            html_lines.append(f'<li>{line[2:]}</li>')
        # Numbered list
        elif re.match(r'^\d+\.\s', line):
            if not in_list:
                html_lines.append('<ul>')
                in_list = True
            content = re.sub(r'^\d+\.\s', '', line)
            html_lines.append(f'<li>{content}</li>')
        # Empty line
        elif not line.strip():
            if in_list:
                html_lines.append('</ul>')
                in_list = False
        # Regular paragraph
        else:
            if in_list:
                html_lines.append('</ul>')
                in_list = False
            if line.strip():
                html_lines.append(f'<p>{line}</p>')
    
    if in_list:
        html_lines.append('</ul>')
    
    return '\n'.join(html_lines)



def export_to_pdf(markdown: str, output_path: str) -> str:
    """
    将 Slidev Markdown 导出为 PDF
    
    使用轻量级方案，适合 512MB 内存限制的环境
    
    Args:
        markdown: Slidev 格式的 Markdown
        output_path: 输出文件路径
    
    Returns:
        生成的 PDF 文件路径
    
    Raises:
        RuntimeError: 如果导出失败
    """
    if not markdown or not markdown.strip():
        raise ValueError("Markdown 内容不能为空")
    
    try:
        # Try using weasyprint (lighter than playwright)
        from weasyprint import HTML
        
        # Render to HTML first
        html_content = render_preview_html(markdown)
        
        # Convert HTML to PDF
        HTML(string=html_content).write_pdf(output_path)
        
        return output_path
        
    except ImportError:
        # Fallback: try reportlab
        try:
            return _export_pdf_reportlab(markdown, output_path)
        except ImportError:
            raise RuntimeError("PDF 导出需要安装 weasyprint 或 reportlab 库")
    except Exception as e:
        raise RuntimeError(f"PDF 导出失败: {str(e)}")


def _export_pdf_reportlab(markdown: str, output_path: str) -> str:
    """
    使用 reportlab 导出 PDF（备用方案）
    """
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.units import inch
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    
    # Parse slides
    slides = parse_slidev_markdown(markdown)
    
    # Create PDF document
    doc = SimpleDocTemplate(
        output_path,
        pagesize=landscape(A4),
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=72
    )
    
    # Get styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=20
    )
    body_style = ParagraphStyle(
        'SlideBody',
        parent=styles['Normal'],
        fontSize=14,
        spaceAfter=10,
        leading=18
    )
    
    # Build story
    story = []
    
    for i, slide in enumerate(slides):
        content = slide.get("content", "")
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line or line.startswith('layout:') or line.startswith('::'):
                continue
            
            if line.startswith('# '):
                story.append(Paragraph(line[2:], title_style))
            elif line.startswith('## '):
                story.append(Paragraph(line[3:], title_style))
            elif line.startswith('- ') or line.startswith('* '):
                story.append(Paragraph(f"• {line[2:]}", body_style))
            elif re.match(r'^\d+\.\s', line):
                story.append(Paragraph(line, body_style))
            else:
                story.append(Paragraph(line, body_style))
        
        # Add page break between slides
        if i < len(slides) - 1:
            story.append(PageBreak())
    
    # Build PDF
    doc.build(story)
    
    return output_path



def export_to_pptx(markdown: str, output_path: str) -> str:
    """
    将 Slidev Markdown 导出为 PPTX
    
    Args:
        markdown: Slidev 格式的 Markdown
        output_path: 输出文件路径
    
    Returns:
        生成的 PPTX 文件路径
    
    Raises:
        RuntimeError: 如果导出失败
    """
    if not markdown or not markdown.strip():
        raise ValueError("Markdown 内容不能为空")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        # Parse slides
        slides = parse_slidev_markdown(markdown)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)
        
        for i, slide_data in enumerate(slides):
            layout = slide_data.get("layout", "default")
            content = slide_data.get("content", "")
            
            # Add slide with blank layout
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Parse content
            lines = content.split('\n')
            title_text = ""
            body_lines = []
            
            for line in lines:
                line = line.strip()
                if not line or line.startswith('layout:') or line.startswith('::'):
                    continue
                
                if line.startswith('# '):
                    title_text = line[2:]
                elif line.startswith('## '):
                    title_text = line[3:]
                else:
                    body_lines.append(line)
            
            # Add title
            if title_text:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5),
                    Inches(12.333), Inches(1.2)
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = title_text
                title_para.font.size = Pt(36)
                title_para.font.bold = True
                if layout == "center":
                    title_para.alignment = PP_ALIGN.CENTER
            
            # Add body content
            if body_lines:
                body_top = Inches(2) if title_text else Inches(0.5)
                
                if layout == "two-cols":
                    # Split content for two columns
                    mid = len(body_lines) // 2
                    left_lines = body_lines[:mid]
                    right_lines = body_lines[mid:]
                    
                    # Left column
                    _add_text_box(slide, left_lines, Inches(0.5), body_top, Inches(5.8), Inches(5))
                    # Right column
                    _add_text_box(slide, right_lines, Inches(6.8), body_top, Inches(5.8), Inches(5))
                else:
                    # Single column
                    _add_text_box(slide, body_lines, Inches(0.5), body_top, Inches(12.333), Inches(5))
        
        # Save presentation
        prs.save(output_path)
        
        return output_path
        
    except ImportError:
        raise RuntimeError("PPTX 导出需要安装 python-pptx 库")
    except Exception as e:
        raise RuntimeError(f"PPTX 导出失败: {str(e)}")


def _add_text_box(slide, lines: list, left, top, width, height):
    """
    向幻灯片添加文本框
    """
    from pptx.util import Pt
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()
        
        # Handle bullet points
        if line.startswith('- ') or line.startswith('* '):
            para.text = f"• {line[2:]}"
            para.level = 0
        elif re.match(r'^\d+\.\s', line):
            para.text = line
            para.level = 0
        else:
            para.text = line
        
        para.font.size = Pt(18)


def count_slides(markdown: str) -> int:
    """
    统计 Slidev Markdown 中的幻灯片数量
    
    Args:
        markdown: Slidev 格式的 Markdown
    
    Returns:
        幻灯片数量
    """
    slides = parse_slidev_markdown(markdown)
    return len(slides)
